"""Opt-in local document extraction, embedding, and cited semantic search."""

from __future__ import annotations

import json
import math
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable, Protocol

from .commands import ActionResult
from .storage import Database, SettingsRepository


SUPPORTED = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx"}


class Embedder(Protocol):
    def embed(self, texts: list[str]) -> list[list[float]]: ...


class OllamaEmbedder:
    def __init__(self, model: str):
        self.model = model

    def embed(self, texts: list[str]) -> list[list[float]]:
        import ollama
        response = ollama.embed(model=self.model, input=texts)
        return response["embeddings"]


class KnowledgeIndex:
    def __init__(self, database: Database, settings: SettingsRepository, embedder: Embedder | None = None):
        self.database = database
        self.settings = settings
        self.embedder = embedder or OllamaEmbedder(settings.get("embedding_model", "nomic-embed-text"))

    def index_configured(self) -> ActionResult:
        roots = [Path(item).expanduser().resolve() for item in self.settings.get("indexed_folders", [])]
        if not roots:
            return ActionResult(False, "No indexed folders are configured in Settings.")
        changed = 0
        errors: list[str] = []
        for root in roots:
            if not root.is_dir():
                errors.append(f"Missing folder: {root}")
                continue
            for path in root.rglob("*"):
                if path.is_file() and path.suffix.lower() in SUPPORTED:
                    try:
                        changed += int(self.index_file(path))
                    except Exception as exc:
                        errors.append(f"{path.name}: {exc}")
        message = f"Indexed {changed} new or changed document(s)."
        if errors:
            message += f" {len(errors)} item(s) could not be indexed."
        return ActionResult(True, message, {"matches": errors[:20]})

    def index_file(self, path: Path) -> bool:
        path = path.resolve()
        modified = path.stat().st_mtime
        with self.database.connect() as connection:
            row = connection.execute("SELECT modified FROM indexed_documents WHERE path=?", (str(path),)).fetchone()
        if row and float(row["modified"]) == modified:
            return False
        extracted = list(self._extract(path))
        chunks: list[tuple[int | None, str]] = []
        for page, content in extracted:
            chunks.extend((page, chunk) for chunk in self._chunks(content))
        if not chunks:
            return False
        embeddings = self.embedder.embed([content for _, content in chunks])
        with self.database.connect() as connection:
            connection.execute("DELETE FROM document_chunks WHERE path=?", (str(path),))
            connection.executemany(
                "INSERT INTO document_chunks(path,page,content,embedding) VALUES(?,?,?,?)",
                [(str(path), page, content, json.dumps(vector)) for (page, content), vector in zip(chunks, embeddings)],
            )
            connection.execute(
                "INSERT INTO indexed_documents(path,modified,indexed_at) VALUES(?,?,?) "
                "ON CONFLICT(path) DO UPDATE SET modified=excluded.modified,indexed_at=excluded.indexed_at",
                (str(path), modified, datetime.now(timezone.utc).isoformat()),
            )
        return True

    def search(self, query: str, limit: int = 8) -> ActionResult:
        query_vector = self.embedder.embed([query])[0]
        with self.database.connect() as connection:
            rows = connection.execute("SELECT path,page,content,embedding FROM document_chunks").fetchall()
        scored = []
        for row in rows:
            score = self._cosine(query_vector, json.loads(row["embedding"]))
            scored.append((score, row))
        scored.sort(key=lambda item: item[0], reverse=True)
        matches = []
        for score, row in scored[:limit]:
            citation = row["path"] + (f"#page={row['page']}" if row["page"] else "")
            excerpt = row["content"].replace("\n", " ")[:280]
            matches.append(f"{citation}\n  {excerpt}  [score {score:.2f}]")
        if not matches:
            return ActionResult(False, "The document index is empty. Add folders in Settings and run indexing.")
        return ActionResult(True, f"Found {len(matches)} relevant document passage(s).", {"matches": matches})

    def clear(self) -> None:
        with self.database.connect() as connection:
            connection.execute("DELETE FROM document_chunks")
            connection.execute("DELETE FROM indexed_documents")

    @staticmethod
    def _chunks(text: str, size: int = 1200, overlap: int = 150) -> Iterable[str]:
        cleaned = " ".join(text.split())
        position = 0
        while position < len(cleaned):
            yield cleaned[position:position + size]
            position += max(1, size - overlap)

    @staticmethod
    def _extract(path: Path) -> Iterable[tuple[int | None, str]]:
        suffix = path.suffix.lower()
        if suffix in {".txt", ".md"}:
            yield None, path.read_text(encoding="utf-8", errors="replace")
        elif suffix == ".pdf":
            from pypdf import PdfReader
            for page_number, page in enumerate(PdfReader(path).pages, start=1):
                yield page_number, page.extract_text() or ""
        elif suffix == ".docx":
            from docx import Document
            yield None, "\n".join(paragraph.text for paragraph in Document(path).paragraphs)
        elif suffix == ".pptx":
            from pptx import Presentation
            for number, slide in enumerate(Presentation(path).slides, start=1):
                yield number, "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            book = load_workbook(path, read_only=True, data_only=True)
            for number, sheet in enumerate(book.worksheets, start=1):
                rows = (" | ".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
                yield number, sheet.title + "\n" + "\n".join(rows)

    @staticmethod
    def _cosine(left: list[float], right: list[float]) -> float:
        numerator = sum(a * b for a, b in zip(left, right))
        denominator = math.sqrt(sum(a * a for a in left)) * math.sqrt(sum(b * b for b in right))
        return numerator / denominator if denominator else 0.0
